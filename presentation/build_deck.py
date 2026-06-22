"""
Build the Panthera AI Agent Classifier presentation deck (.pptx).

15-minute professional presentation for the paper's authors + internship tutor.
Top-down narrative: Goal -> Methodology -> Platform -> Results -> Reflections
-> Weaknesses -> Conclusions.

Design mirrors the live platform: dark UI (#1A1A1A) with copper accent (#C4845A)
and the advantage palette (informational/analytical/behavioral).

Run:  python presentation/build_deck.py
Out:  presentation/Panthera_AI_Agent_Classifier.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ── Paths ────────────────────────────────────────────────────────────────
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")
IMG = os.path.join(ROOT, "ai_agent_classifier", "static", "img")
OUT = os.path.join(ROOT, "presentation", "Panthera_AI_Agent_Classifier.pptx")

# ── Palette (from static/css/style.css) ──────────────────────────────────
BG        = RGBColor(0x1A, 0x1A, 0x1A)
SURFACE   = RGBColor(0x23, 0x23, 0x23)
ELEVATED  = RGBColor(0x2A, 0x2A, 0x2A)
BORDER    = RGBColor(0x36, 0x36, 0x36)
ACCENT    = RGBColor(0xC4, 0x84, 0x5A)   # copper
ACCENT_HI = RGBColor(0xD6, 0x98, 0x72)
ACCENT_DK = RGBColor(0x8B, 0x5E, 0x3C)
TXT       = RGBColor(0xF0, 0xF0, 0xF0)
TXT2      = RGBColor(0xA0, 0xA0, 0xA0)
TXT_MUTE  = RGBColor(0x66, 0x66, 0x66)
INF       = RGBColor(0xEF, 0x9F, 0x27)   # informational - orange
ANA       = RGBColor(0x1D, 0xB8, 0x7A)   # analytical - green
BEH       = RGBColor(0xE0, 0x61, 0x3A)   # behavioral - rust
NEARBLACK = RGBColor(0x11, 0x11, 0x11)

FONT = "Segoe UI"        # platform's fallback; present on every Windows box
FONT_LT = "Segoe UI Light"

# ── Geometry (16:9) ──────────────────────────────────────────────────────
EMU_IN = 914400
SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _no_autosize(tf):
    tf.word_wrap = True
    el = tf._txBody
    bodyPr = el.find(qn('a:bodyPr'))
    for tag in ('a:spAutoFit', 'a:normAutofit'):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(el.makeelement(qn('a:noAutofit'), {}))


def box(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    _no_autosize(tf)
    return tb, tf


def run(p, text, size, color=TXT, bold=False, font=FONT, italic=False, spacing=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set('spc', str(int(spacing * 100)))
    return r


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def rrect(s, l, t, w, h, fill=None, line=None, line_w=1.0, radius=0.08):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def eyebrow(s, text, top=MARGIN, left=MARGIN, color=ACCENT):
    _, tf = box(s, left, top, Inches(9), Inches(0.3))
    run(tf.paragraphs[0], text.upper(), 12, color, bold=True, spacing=2.2)


def title(s, text, top=Inches(1.02), left=MARGIN, size=30, width=Inches(11.6)):
    _, tf = box(s, left, top, width, Inches(1.0))
    run(tf.paragraphs[0], text, size, TXT, bold=True)
    # copper underline accent
    rect(s, left + Emu(2000), top + Inches(0.82), Inches(0.62), Pt(3.2), fill=ACCENT)


def footer(s, n):
    rect(s, 0, SH - Inches(0.34), SW, Pt(2.2), fill=ACCENT_DK)
    _, tf = box(s, MARGIN, SH - Inches(0.30), Inches(7), Inches(0.26))
    run(tf.paragraphs[0], "PANTHERA  ·  Multi-Dimensional AI Agent Classifier", 9, TXT_MUTE, spacing=1.4)
    _, tf2 = box(s, SW - Inches(2.2), SH - Inches(0.30), Inches(1.58), Inches(0.26))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run(p, f"{n:02d} / 15", 9, TXT_MUTE, spacing=1.4)


def bullets(tf, items, size=15, gap=10, color=TXT, lead=ACCENT):
    """items: list of (text, level) or (head, body) tuples handled by caller."""
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap); p.space_before = Pt(0)
        if isinstance(it, tuple):
            head, body = it
            run(p, "▸  ", size, lead, bold=True)
            run(p, head, size, color, bold=True)
            if body:
                run(p, body, size, TXT2)
        else:
            run(p, "▸  ", size, lead, bold=True)
            run(p, it, size, color)
    return tf


def fit_image(s, path, l, t, w, h, frame=True):
    """Place image scaled to fit inside (l,t,w,h), centered, with subtle frame."""
    iw, ih = Image.open(path).size
    boxr = w / h
    imgr = iw / ih
    if imgr > boxr:
        nw = w; nh = int(w / imgr)
    else:
        nh = h; nw = int(h * imgr)
    nl = l + (w - nw) // 2
    nt = t + (h - nh) // 2
    if frame:
        rect(s, nl - Emu(9000), nt - Emu(9000), nw + Emu(18000), nh + Emu(18000),
             fill=NEARBLACK, line=ACCENT_DK, line_w=1.25)
    s.shapes.add_picture(path, nl, nt, nw, nh)
    return nl, nt, nw, nh


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def chip(s, l, t, label, color, w=Inches(1.9), h=Inches(0.4)):
    c = rrect(s, l, t, w, h, fill=None, line=color, line_w=1.5, radius=0.5)
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, label, 12, color, bold=True)
    return c


def statcard(s, l, t, w, h, value, label, vcolor=ACCENT):
    rrect(s, l, t, w, h, fill=SURFACE, line=BORDER, line_w=1.0, radius=0.06)
    _, tf = box(s, l + Inches(0.18), t + Inches(0.16), w - Inches(0.36),
                h - Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.space_after = Pt(2)
    run(p, value, 30, vcolor, bold=True)
    p2 = tf.add_paragraph()
    run(p2, label, 11, TXT2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════
s = slide()
# top identity strip
rect(s, 0, 0, SW, Inches(0.07), fill=ACCENT)
logo = os.path.join(IMG, "panthera_logo.png")
if os.path.exists(logo):
    lw = Inches(2.5); lh = Inches(2.5 * 610 / 1562)
    s.shapes.add_picture(logo, MARGIN, Inches(0.66), lw, lh)

_, tf = box(s, MARGIN, Inches(2.55), Inches(11.8), Inches(0.4))
run(tf.paragraphs[0], "RESEARCH FRAMEWORK · OPERATIONALIZED", 13, ACCENT, bold=True, spacing=2.5)

_, tf = box(s, MARGIN, Inches(3.05), Inches(11.9), Inches(2.0))
p = tf.paragraphs[0]; p.space_after = Pt(4)
run(p, "A Multi-Dimensional Classification System", 40, TXT, bold=True)
p2 = tf.add_paragraph()
run(p2, "for AI Agents in the Investment Industry", 40, ACCENT, bold=True)

_, tf = box(s, MARGIN, Inches(5.05), Inches(11.9), Inches(0.7))
run(tf.paragraphs[0],
    "From the Schuller · Wierckx · Kuhn · Zilic (2025) framework to a live, "
    "queryable benchmark of 78 AI agents.", 16, TXT2)

rect(s, MARGIN, Inches(5.95), Inches(3.2), Pt(2), fill=ACCENT_DK)
_, tf = box(s, MARGIN, Inches(6.1), Inches(11.9), Inches(0.8))
p = tf.paragraphs[0]
run(p, "Andrea Signoretti", 14, TXT, bold=True)
run(p, "   ·   Panthera Solutions Internship", 14, TXT2)
p2 = tf.add_paragraph(); p2.space_before = Pt(2)
run(p2, "Presented to the authors of the framework & internship tutor   ·   June 2026", 12, TXT_MUTE)
notes(s,
"Good [morning/afternoon], and thank you for the time. My name is Andrea Signoretti. "
"Over my internship I took your paper — the multi-dimensional classification system for AI agents — "
"and turned it from a conceptual framework into a working software platform that classifies, visualizes "
"and benchmarks the real AI-agent market in investment management. "
"In the next 15 minutes I'll walk top-down: first the goal, then the method, then the platform itself, "
"then what the data actually shows us about the market, my own design decisions, the honest limitations, "
"and where this can go next. I built this directly on your four dimensions, so my hope is this is the "
"empirical validation step your Discussion section calls for.")
footer(s, 1)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE GOAL
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Why this project")
title(s, "The goal: turn a taxonomy into a living instrument")

_, tf = box(s, MARGIN, Inches(2.05), Inches(5.55), Inches(4.4))
bullets(tf, [
    ("Framework in action.  ", "Take the paper's dimensions out of the page and make them "
     "executable — every agent gets a concrete, comparable coordinate."),
    ("A rich benchmark.  ", "One catalogue of the AI-agent market for investment management — "
     "commercial, in-house and academic — classified consistently."),
    ("A market view.  ", "Reveal where AI actually clusters across the investment process, and "
     "where it is conspicuously absent."),
    ("A shared language.  ", "Give managers, developers and regulators one vocabulary to compare "
     "tools and align adoption with a firm's real edge."),
], size=15, gap=14)

# right: target users card
cx = Inches(6.55); cw = Inches(6.2)
rrect(s, cx, Inches(2.05), cw, Inches(4.3), fill=SURFACE, line=BORDER, radius=0.05)
_, tf = box(s, cx + Inches(0.4), Inches(2.4), cw - Inches(0.8), Inches(0.5))
run(tf.paragraphs[0], "ONE INSTRUMENT, THREE AUDIENCES", 12, ACCENT, bold=True, spacing=1.8)
rows = [
    ("Investment managers", "Which tool fits which decision — and does it match our edge?", INF),
    ("Developers / vendors", "Where is the market crowded, where is the white space?", ANA),
    ("Regulators", "How autonomous is this system, and where does it touch compliance?", BEH),
]
yy = Inches(3.05)
for head, body, col in rows:
    rect(s, cx + Inches(0.4), yy + Inches(0.05), Pt(3.5), Inches(0.9), fill=col)
    _, tf = box(s, cx + Inches(0.62), yy, cw - Inches(1.05), Inches(1.0))
    p = tf.paragraphs[0]; p.space_after = Pt(2)
    run(p, head, 15, TXT, bold=True)
    p2 = tf.add_paragraph()
    run(p2, body, 12.5, TXT2)
    yy += Inches(1.07)
notes(s,
"The starting point was a gap your paper names directly: AI tools are flooding into investment workflows, "
"but there is no consistent way to compare or categorize them. So the goal was deliberately practical. "
"Four things. One — get the framework 'in action': make each of your dimensions something you can actually "
"assign and compute, not just describe. Two — build a rich benchmark: a single, consistently-classified "
"catalogue of the market. Three — from that catalogue, surface a market view: where AI concentrates and "
"where it's missing. Four — deliver the shared language your paper argues for, usable by the three audiences "
"on the right: managers choosing tools, developers and vendors hunting white space, and regulators assessing "
"autonomy and compliance exposure. Everything that follows serves these four aims.")
footer(s, 2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The problem")
title(s, "Without a shared taxonomy, firms misjudge AI")
_, tf = box(s, MARGIN, Inches(1.95), Inches(11.8), Inches(0.6))
run(tf.paragraphs[0],
    "The paper's core claim: absent a common structure, firms underutilize AI or misjudge its "
    "strategic potential. Three failure modes motivate the build.", 15, TXT2)

cards = [
    ("Wrong tool, wrong context",
     "A white-swan optimizer deployed against a black-swan problem — confident answers where no "
     "distribution exists.", BEH),
    ("Over-automation",
     "Handing genuine judgment calls to systems built for routine, well-characterized tasks.", INF),
    ("Under-utilization",
     "Leaving durable edge on the table where AI could reliably compound an advantage.", ANA),
]
cw = Inches(3.86); gap = Inches(0.2); cx = MARGIN; cy = Inches(2.75); ch = Inches(2.9)
for head, body, col in cards:
    rrect(s, cx, cy, cw, ch, fill=SURFACE, line=BORDER, radius=0.05)
    rect(s, cx, cy, cw, Pt(4), fill=col)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(0.4), cw - Inches(0.6), Inches(0.7))
    run(tf.paragraphs[0], head, 17, TXT, bold=True)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(1.25), cw - Inches(0.6), Inches(1.5))
    run(tf.paragraphs[0], body, 13.5, TXT2)
    cx += cw + gap

_, tf = box(s, MARGIN, Inches(5.95), Inches(11.8), Inches(0.8))
p = tf.paragraphs[0]
run(p, "My read:  ", 14, ACCENT, bold=True)
run(p, "each failure is a classification error before it is an investment error. "
       "If you can place a tool precisely, you can govern it. The platform makes that placement cheap "
       "and repeatable.", 14, TXT)
notes(s,
"Why does an instrument like this matter? Because the paper frames the risk as strategic, and I'd push it "
"one step further. Look at the three failure modes. Wrong tool for the context — a white-swan optimizer, "
"which assumes a knowable distribution, pointed at a black-swan problem where no distribution exists. "
"Over-automation — giving real judgment to a tool designed for routine tasks. And under-utilization — the "
"quiet one — leaving durable edge unused. My interpretation is the line at the bottom: every one of these is "
"a classification error before it becomes an investment error. Misplace the tool on the map and the bad "
"outcome follows. So if placing a tool precisely becomes cheap and repeatable, governance becomes possible. "
"That's the thesis the platform is built to serve.")
footer(s, 3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE FRAMEWORK
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The framework I operationalized")
title(s, "Five dimensions — your paper, made assignable")

dims = [
    ("COMPLEXITY — Swan Theory", ACCENT,
     "White · Light-grey · Dark-grey · Black", "How knowable is the outcome distribution? (Knight · Taleb)"),
    ("COMPARATIVE ADVANTAGE", ANA,
     "Informational · Analytical · Behavioral", "What edge does the agent confer? One primary edge only."),
    ("AUTONOMY", INF,
     "Low · Medium · High · Full", "How much human veto remains over the investment action?"),
    ("INVESTMENT PROCESS STAGE", BEH,
     "7 stages: idea-gen → stakeholder", "Where in the lifecycle does it intervene? Multi-select."),
    ("AGENT TYPE  (provenance)", TXT2,
     "Commercial · In-house · Academic", "Origin of the system, for benchmark integrity."),
]
yy = Inches(2.0); rh = Inches(0.92)
for name, col, vals, desc in dims:
    rrect(s, MARGIN, yy, Inches(11.95), rh - Inches(0.12), fill=SURFACE, line=BORDER, radius=0.07)
    rect(s, MARGIN, yy, Pt(4.5), rh - Inches(0.12), fill=col)
    _, tf = box(s, MARGIN + Inches(0.28), yy + Inches(0.1), Inches(3.6), rh - Inches(0.3),
                anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], name, 13.5, col, bold=True, spacing=0.6)
    _, tf = box(s, MARGIN + Inches(3.95), yy + Inches(0.1), Inches(3.4), rh - Inches(0.3),
                anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], vals, 13, TXT, bold=True)
    _, tf = box(s, MARGIN + Inches(7.45), yy + Inches(0.1), Inches(4.3), rh - Inches(0.3),
                anchor=MSO_ANCHOR.MIDDLE)
    run(tf.paragraphs[0], desc, 12, TXT2)
    yy += rh
notes(s,
"Here are the dimensions the platform makes assignable — four orthogonal ones plus provenance. "
"Complexity uses your Swan Theory: the key question isn't whether uncertainty exists, but whether the "
"probability distribution of outcomes can be characterized — white is classical risk, black is radical "
"uncertainty, grounded in Knight and Taleb. Comparative advantage — informational, analytical, behavioral — "
"the Fuller and Wierckx typology, and I enforce a single primary edge, which I'll defend later. Autonomy is "
"about human veto over the investment action. Stages are the seven-step pipeline, multi-select because real "
"tools span several. And agent type tracks origin so the benchmark stays honest about what's a shipping "
"product versus a research prototype. The whole platform is just these five fields, applied rigorously, "
"77 times over.")
footer(s, 4)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Methodology")
title(s, "How an agent becomes a data point")

steps = [
    ("1", "SOURCE", "Scan the market — vendor products, bank in-house systems, academic prototypes — "
     "and capture name + URL.", ACCENT),
    ("2", "AI-ASSISTED DRAFT", "Claude (claude-opus-4-8) researches each agent and returns a structured "
     "classification via forced tool-use — never free text.", INF),
    ("3", "HUMAN REVIEW", "I verify every call against the framework guide, focusing on the nuanced "
     "complexity tier and the rationale.", ANA),
    ("4", "PERSIST & PUBLISH", "Stored in SQLite; only 'classified' agents surface in the matrix, finder "
     "and exports.", BEH),
]
cw = Inches(2.92); gap = Inches(0.12); cx = MARGIN; cy = Inches(2.05); ch = Inches(3.05)
for num, head, body, col in steps:
    rrect(s, cx, cy, cw, ch, fill=SURFACE, line=BORDER, radius=0.05)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.3), cy + Inches(0.3),
                              Inches(0.62), Inches(0.62))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    circ.shadow.inherit = False
    ctf = circ.text_frame; ctf.word_wrap = False
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    run(cp, num, 20, NEARBLACK, bold=True)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(0.5))
    run(tf.paragraphs[0], head, 14, col, bold=True, spacing=0.8)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(1.55), cw - Inches(0.6), Inches(1.4))
    run(tf.paragraphs[0], body, 12.5, TXT2)
    if cx + cw < Inches(12):
        ar = box(s, cx + cw - Inches(0.02), cy + Inches(1.2), Inches(0.2), Inches(0.4))[1]
        run(ar.paragraphs[0], "›", 22, ACCENT_DK, bold=True)
    cx += cw + gap

_, tf = box(s, MARGIN, Inches(5.45), Inches(11.9), Inches(1.2))
p = tf.paragraphs[0]
run(p, "Design principle:  ", 14, ACCENT, bold=True)
run(p, "AI for scale, human for judgment.  ", 14, TXT, bold=True)
run(p, "The model does the breadth — researching dozens of tools — but the framework is mine to "
       "uphold. Forced tool-use guarantees a comparable shape for every record; the human pass guards "
       "the meaning. This is augmented intelligence applied to building the benchmark itself.", 14, TXT2)
notes(s,
"This is the pipeline that turns a product name into a comparable data point. Four steps. I source agents "
"across the three provenance types. Then — and this is the part I want to be transparent about — I use Claude "
"to draft the classification. It researches each agent and must answer through a structured tool call, so "
"every record comes back in the same shape, with the same fields. It cannot ramble in free text. Step three, "
"I review every single one against the classification guide, paying most attention to the complexity tier, "
"which is the subtlest dimension, and to the rationale. Then it's persisted, and only fully-classified agents "
"are published to the views. The principle at the bottom is the honest summary: AI for scale, human for "
"judgment. There's a nice reflexivity here — the benchmark about augmented intelligence is itself built with "
"augmented intelligence. I'll come back to the limitations of leaning on the model in step two.")
footer(s, 5)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE PLATFORM (overview)
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The platform")
title(s, "Four lenses on one catalogue")
fit_image(s, os.path.join(SHOTS, "framework.png"),
          Inches(6.5), Inches(1.95), Inches(6.3), Inches(4.55))

_, tf = box(s, MARGIN, Inches(2.05), Inches(5.5), Inches(4.4))
bullets(tf, [
    ("Matrix.  ", "The multi-dimensional classification grid — complexity × stage, every agent placed."),
    ("Guide / Agent Finder.  ", "Reverse lookup — describe a need, get the right tools."),
    ("Framework.  ", "The living methodology page, citing the source paper."),
    ("Exports.  ", "One-click Excel, Word and PDF reports for stakeholders."),
], size=15, gap=12)
_, tf = box(s, MARGIN, Inches(5.25), Inches(5.5), Inches(1.3))
p = tf.paragraphs[0]
run(p, "Stack:  ", 13, ACCENT, bold=True)
run(p, "Flask · SQLAlchemy · SQLite · Bootstrap 5 · Anthropic Claude API. "
       "Self-hostable, no external dependencies at runtime beyond the browser.", 13, TXT2)
notes(s,
"Here's the platform itself. It's a self-hostable Flask web app, and it gives you four lenses onto the same "
"catalogue of agents. The Matrix is the flagship — the complexity-by-stage grid. The Guide, or Agent Finder, "
"runs the framework in reverse: instead of 'what is this tool', you ask 'what do I need' and it filters down. "
"The Framework page keeps the methodology live and cites your paper directly, so the provenance is always one "
"click away. And Exports turn the whole thing into board-ready Excel, Word and PDF. The stack is deliberately "
"boring and robust — Flask, SQLite, Bootstrap, plus the Claude API for the classification step. Nothing exotic "
"to run. Let me show you each lens.")
footer(s, 6)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MATRIX
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The platform · 1")
title(s, "The Multi-Dimensional Matrix — three signals in one badge")
fit_image(s, os.path.join(SHOTS, "matrix.png"),
          MARGIN, Inches(1.95), Inches(8.1), Inches(4.55))

px = Inches(9.0); pw = Inches(3.8)
_, tf = box(s, px, Inches(2.0), pw, Inches(0.5))
run(tf.paragraphs[0], "EVERY BADGE ENCODES:", 12, ACCENT, bold=True, spacing=1.5)
enc = [
    ("Fill colour", "Comparative advantage", ANA),
    ("Border style", "Agent type (solid/outline/dashed)", INF),
    ("Corner pip", "Autonomy level", BEH),
]
yy = Inches(2.5)
for head, body, col in enc:
    rect(s, px, yy + Inches(0.05), Pt(3.5), Inches(0.62), fill=col)
    _, tf = box(s, px + Inches(0.2), yy, pw - Inches(0.2), Inches(0.7))
    p = tf.paragraphs[0]; p.space_after = Pt(1)
    run(p, head, 13.5, TXT, bold=True)
    run(tf.add_paragraph(), body, 12, TXT2)
    yy += Inches(0.78)

rrect(s, px, Inches(5.0), pw, Inches(1.5), fill=SURFACE, line=BORDER, radius=0.06)
_, tf = box(s, px + Inches(0.22), Inches(5.18), pw - Inches(0.44), Inches(1.2))
p = tf.paragraphs[0]
run(p, "Rows = ", 12.5, TXT2); run(p, "Swan complexity", 12.5, TXT, bold=True)
run(p, ".  Columns = ", 12.5, TXT2); run(p, "7 process stages", 12.5, TXT, bold=True)
run(p, ".  One screen reads the whole market at a glance.", 12.5, TXT2)
notes(s,
"This is the Matrix, and it's the densest part of the whole project, so let me decode it. Rows are the four "
"Swan complexity tiers. Columns are the seven investment-process stages. So position alone tells you 'how "
"hard is this problem, and where in the workflow does this tool act'. But I didn't want to lose the other "
"dimensions, so each agent badge carries three more signals simultaneously: the fill colour is its "
"comparative advantage, the border style is its provenance — solid commercial, outline in-house, dashed "
"academic — and a small pip in the corner encodes autonomy, with a pulsing crimson pip for the fully "
"autonomous systems. That's five dimensions legible in a single chip, and the entire market on one screen. "
"This information-density was a deliberate choice and I'll own it in the reflections — it's powerful but it "
"has a learning curve.")
footer(s, 7)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ADVANTAGE & AUTONOMY VIEWS
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The platform · 2")
title(s, "Re-pivot the same data: advantage & autonomy")
fit_image(s, os.path.join(SHOTS, "advantage.png"),
          MARGIN, Inches(2.0), Inches(6.05), Inches(3.0))
fit_image(s, os.path.join(SHOTS, "autonomy.png"),
          Inches(6.85), Inches(2.0), Inches(5.95), Inches(3.0))
_, tf = box(s, MARGIN, Inches(5.15), Inches(6.05), Inches(1.3))
p = tf.paragraphs[0]; p.space_after = Pt(3)
run(p, "Advantage view", 14, ANA, bold=True)
run(tf.add_paragraph(),
    "Stage × advantage — which kind of edge dominates each phase of the process.", 13, TXT2)
_, tf = box(s, Inches(6.85), Inches(5.15), Inches(5.95), Inches(1.3))
p = tf.paragraphs[0]; p.space_after = Pt(3)
run(p, "Autonomy scatter", 14, BEH, bold=True)
run(tf.add_paragraph(),
    "Agents ranked by how much human veto they surrender — fully autonomous systems isolated at the top.",
    13, TXT2)
notes(s,
"A single arrangement can hide as much as it shows, so the same catalogue re-pivots. On the left, the "
"advantage view holds the stages as columns but swaps the rows to the three edge types, so you can read "
"which kind of advantage dominates each phase — for instance, informational tools concentrating at idea "
"generation. On the right, the autonomy scatter ranks every agent by how much human veto it gives up, and "
"deliberately isolates the fully-autonomous systems — the AIEQs and Numerais of the world — in their own "
"band at the top, because those are the ones a regulator or a CIO needs to find first. Same data, three "
"questions answered. This is the framework being genuinely multi-dimensional rather than a single fixed chart.")
footer(s, 8)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GUIDE / AGENT FINDER
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The platform · 3")
title(s, "Agent Finder — the framework, run backwards")
fit_image(s, os.path.join(SHOTS, "guide.png"),
          Inches(6.4), Inches(1.95), Inches(6.4), Inches(4.5))
_, tf = box(s, MARGIN, Inches(2.05), Inches(5.6), Inches(4.5))
bullets(tf, [
    ("From description to shortlist.  ", "Pick a stage, an edge, an autonomy ceiling and a complexity "
     "tier — get a filtered, ranked shortlist instantly."),
    ("Decision support, not a catalogue.  ", "This is where the benchmark earns its keep: it answers "
     "\"what should we use here?\""),
    ("Governance-aware.  ", "Cap autonomy or require a compliance-stage fit before a tool even appears."),
], size=15, gap=14)
_, tf = box(s, MARGIN, Inches(5.5), Inches(5.6), Inches(1.0))
p = tf.paragraphs[0]
run(p, "My reasoning:  ", 13.5, ACCENT, bold=True)
run(p, "classification is only half the value. The other half is retrieval — turning the map into a "
       "recommendation a manager can act on.", 13.5, TXT2)
notes(s,
"The Matrix answers 'what is this tool'. The Finder answers the question a working manager actually asks: "
"'I have this need — what should I use?' You select the stage you're in, the kind of edge you want, an "
"autonomy ceiling, a complexity tier, and it returns a ranked shortlist immediately. Crucially it's "
"governance-aware: you can cap autonomy, or require that a tool fits the compliance stage, before it's even "
"eligible to appear. My reasoning here was simple — classification is only half the value. A taxonomy that "
"just sits there is an encyclopedia. The retrieval step is what turns the map into a decision a manager can "
"defend. That's the difference between a catalogue and decision support.")
footer(s, 9)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EXPORTS
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The platform · 4")
title(s, "From screen to stakeholder in one click")
fit_image(s, os.path.join(SHOTS, "pdf_export.png"),
          Inches(8.7), Inches(1.95), Inches(4.1), Inches(4.5))
_, tf = box(s, MARGIN, Inches(2.05), Inches(7.7), Inches(0.55))
run(tf.paragraphs[0],
    "The framework's stakeholder-management stage applied to itself — auditable, shareable artefacts.",
    14, TXT2)

exps = [
    ("Excel  (.xlsx)", "Three sheets — colour-coded matrix, tabular detail, summary statistics. "
     "Frozen panes, hyperlinks.", ANA),
    ("Word  (.docx)", "Cover, table of contents, one section per agent: classification table, key "
     "features, stage timeline.", INF),
    ("PDF  (.pdf)", "Same structure, ReportLab-rendered with internal links — the board-ready handout.",
     BEH),
]
yy = Inches(2.85)
for head, body, col in exps:
    rrect(s, MARGIN, yy, Inches(7.7), Inches(1.05), fill=SURFACE, line=BORDER, radius=0.06)
    rect(s, MARGIN, yy, Pt(4), Inches(1.05), fill=col)
    _, tf = box(s, MARGIN + Inches(0.3), yy + Inches(0.13), Inches(7.2), Inches(0.4))
    run(tf.paragraphs[0], head, 15, TXT, bold=True)
    _, tf = box(s, MARGIN + Inches(0.3), yy + Inches(0.52), Inches(7.2), Inches(0.5))
    run(tf.paragraphs[0], body, 12.5, TXT2)
    yy += Inches(1.18)
notes(s,
"Last lens: reporting. A live web app is great for exploration, but a CIO or a regulator wants an artefact "
"they can file and audit. So with one click the catalogue exports to Excel, Word or PDF. The Excel has the "
"colour-coded matrix, the full tabular detail, and a statistics sheet. The Word and PDF are proper documents "
"— cover, contents, a section per agent with its classification table and rationale. There's a small "
"conceptual joke I like here: this is the framework's own stakeholder-management stage, applied to the "
"framework. Only fully-classified agents are ever exported, so the reports stay clean. This is what makes the "
"tool usable by the regulator audience the paper keeps pointing to.")
footer(s, 10)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RESULTS (the market view)
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Results · what the benchmark reveals")
title(s, "A map of where AI actually lives in investing")

# top stat row
sw_ = Inches(2.86); gp = Inches(0.13); sx = MARGIN; sy = Inches(1.95)
statcard(s, sx, sy, sw_, Inches(1.15), "78", "agents classified", ACCENT); sx += sw_ + gp
statcard(s, sx, sy, sw_, Inches(1.15), "~99%", "White / Light-grey swan", INF); sx += sw_ + gp
statcard(s, sx, sy, sw_, Inches(1.15), "63%", "Compete on analytical edge", ANA); sx += sw_ + gp
statcard(s, sx, sy, sw_, Inches(1.15), "69%", "Low autonomy (copilots)", BEH)

_, tf = box(s, MARGIN, Inches(3.4), Inches(6.05), Inches(3.0))
bullets(tf, [
    ("AI hugs the knowable.  ", "39 white + 40 light-grey, but only 1 dark-grey and 0 black. AI "
     "concentrates exactly where a distribution exists."),
    ("Behavioral edge is a desert.  ", "49 analytical, 30 informational — and just 1 behavioral. "
     "The most durable edge is the least automated."),
    ("Copilots, not autopilots.  ", "54 low + 14 medium autonomy; only 2 fully autonomous. The market "
     "keeps humans in the loop."),
    ("The funnel is front-loaded.  ", "Idea-assessment (56) and idea-generation (31) dominate; "
     "execution (13) trails."),
], size=14, gap=11)

# right interpretation panel
px = Inches(6.95); pw = Inches(5.8)
rrect(s, px, Inches(3.4), pw, Inches(2.95), fill=SURFACE, line=ACCENT_DK, radius=0.05)
_, tf = box(s, px + Inches(0.32), Inches(3.62), pw - Inches(0.64), Inches(0.5))
run(tf.paragraphs[0], "MY INTERPRETATION", 12, ACCENT, bold=True, spacing=1.8)
_, tf = box(s, px + Inches(0.32), Inches(4.08), pw - Inches(0.64), Inches(2.2))
p = tf.paragraphs[0]; p.space_after = Pt(8)
run(p, "The data confirms the paper's thesis empirically: ", 14, TXT)
run(p, "AI's impact is greatest where standardization is high and the distribution is knowable — "
       "and thins out toward genuine uncertainty.", 14, TXT, bold=True)
p2 = tf.add_paragraph()
run(p2, "The white space is the behavioral / black-swan corner. That is precisely where a human edge "
        "— and Panthera's Augmented-Intelligence thesis — still matters most.", 13.5, TXT2)
notes(s,
"Now the payoff — what 78 classified agents actually tell us about the market. Four findings. First, AI hugs "
"the knowable: 39 white-swan and 40 light-grey agents, but only one dark-grey and zero black. Tools cluster "
"exactly where a probability distribution exists. Second, behavioral edge is a desert — 49 analytical, 30 "
"informational, and a single behavioral agent. Your paper argues the behavioral advantage is the most "
"durable; the market shows it's also the least automated. Third, copilots not autopilots — nearly 70% are "
"low-autonomy; only two are fully autonomous. Fourth, the activity is front-loaded into idea assessment and "
"generation, with execution trailing. And the interpretation on the right is the one I most want to leave you "
"with: this is empirical confirmation of your central claim. AI's impact concentrates where standardization "
"is high and the distribution is knowable, and it thins toward real uncertainty. Which means the white space "
"— the behavioral, black-swan corner — is exactly where the human edge, and Panthera's augmented-intelligence "
"thesis, still matters most. The benchmark didn't just catalogue the market; it located the frontier.")
footer(s, 11)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PERSONAL REFLECTIONS / DESIGN DECISIONS
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Reflections · the choices I made")
title(s, "Where I extended the framework — and why")
decisions = [
    ("I added a 4th autonomy tier: 'Full'",
     "The paper implies low→high. I split off systems with zero human veto on live capital (AIEQ, "
     "Numerai). Governance treats 'acts within my limits' and 'sets its own limits' very differently.", BEH),
    ("One advantage per agent — enforced",
     "I refused to let a tool be 'a bit of everything'. Forcing the single primary edge is what makes "
     "the benchmark comparable and the analyst think.", ANA),
    ("Five dimensions in one badge",
     "A density bet: colour, border and pip on every chip. It rewards a trained eye and pays off on the "
     "market-wide view — at the cost of a learning curve.", INF),
    ("AI-drafted, human-ratified",
     "Using Claude to scale classification, but never to finalize it. The framework's authority stays "
     "human; the model only does the legwork.", ACCENT),
]
cw = Inches(6.07); ch = Inches(1.95); gx = Inches(0.2); gy = Inches(0.18)
positions = [(MARGIN, Inches(2.0)), (MARGIN + cw + gx, Inches(2.0)),
             (MARGIN, Inches(2.0) + ch + gy), (MARGIN + cw + gx, Inches(2.0) + ch + gy)]
for (head, body, col), (cx, cy) in zip(decisions, positions):
    rrect(s, cx, cy, cw, ch, fill=SURFACE, line=BORDER, radius=0.05)
    rect(s, cx, cy, Pt(4.5), ch, fill=col)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(0.22), cw - Inches(0.6), Inches(0.55))
    run(tf.paragraphs[0], head, 15.5, TXT, bold=True)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(0.82), cw - Inches(0.6), Inches(1.05))
    run(tf.paragraphs[0], body, 12.5, TXT2)
notes(s,
"I want to be explicit about where I made judgment calls, because these are mine to defend. Four of them. "
"One — I added a fourth autonomy tier, 'Full'. Your text implies a low-to-high range; I split off a top tier "
"for systems with no human veto on live capital, like AIEQ and Numerai. My reasoning: from a governance and "
"regulatory standpoint, 'acts within limits I set' and 'sets its own limits' are categorically different "
"animals, and collapsing them loses the thing a regulator cares about most. Two — I enforce exactly one "
"comparative advantage per agent. It would've been easier to let tools be a bit of everything, but that "
"destroys comparability and lets the analyst off the hook. Forcing the primary edge is a feature. Three — the "
"five-signals-in-one-badge density. That's a genuine bet, and I'll concede the cost: a learning curve. But on "
"a market-wide view it's what lets you see structure instead of scrolling. Four — AI drafts, human ratifies. "
"I used the model for reach, never for the final word. The framework's authority stays human. You may "
"disagree with any of these, and I'd genuinely like to hear it.")
footer(s, 12)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — WEAKNESSES / LIMITATIONS
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Limitations · an honest method critique")
title(s, "What this method does not yet prove")
lims = [
    ("Classifier subjectivity", "The complexity tier especially is a judgment call. Two analysts could "
     "disagree — there is no inter-rater reliability score yet."),
    ("Single-source drafting", "AI-assisted classification inherits the model's knowledge cutoff and "
     "blind spots; sparse-doc agents risk under-classification."),
    ("Self-reported, unvalidated", "Classifications reflect vendor and paper claims, not audited "
     "behaviour. No agent was independently benchmarked on live performance."),
    ("A snapshot, not a film", "78 agents at one moment. The market moves weekly; autonomy and "
     "capability drift. There is no longitudinal dimension yet."),
    ("Thin tails", "Only 1 behavioral and ~0 black-swan agents — too few to generalize about the most "
     "important corner of the map."),
    ("Forced single advantage", "The very choice that aids comparability also flattens genuinely "
     "hybrid tools. Comparability bought with nuance."),
]
cw = Inches(3.93); ch = Inches(1.5); gx = Inches(0.13); gy = Inches(0.16)
sx0 = MARGIN; sy0 = Inches(2.0)
for i, (head, body) in enumerate(lims):
    cx = sx0 + (i % 3) * (cw + gx)
    cy = sy0 + (i // 3) * (ch + gy)
    rrect(s, cx, cy, cw, ch, fill=SURFACE, line=BORDER, radius=0.05)
    _, tf = box(s, cx + Inches(0.26), cy + Inches(0.18), cw - Inches(0.5), Inches(0.4))
    p = tf.paragraphs[0]
    run(p, "!  ", 14, BEH, bold=True)
    run(p, head, 14.5, TXT, bold=True)
    _, tf = box(s, cx + Inches(0.26), cy + Inches(0.62), cw - Inches(0.5), Inches(0.82))
    run(tf.paragraphs[0], body, 12, TXT2)

_, tf = box(s, MARGIN, Inches(5.55), Inches(11.9), Inches(0.9))
p = tf.paragraphs[0]
run(p, "Stance:  ", 14, ACCENT, bold=True)
run(p, "none of these invalidate the instrument — they scope it. It is a structured, transparent "
       "first map, explicitly built to be challenged and refined. That is the validation work your paper "
       "calls for.", 13.5, TXT)
notes(s,
"I'd rather name the weaknesses than have you find them. Six. The classifier is subjective — the complexity "
"tier especially is a judgment call, and I don't yet have an inter-rater reliability score. The drafting leans "
"on a single model, so it inherits that model's cutoff and blind spots, and thinly-documented agents can get "
"under-classified. The classifications are self-reported — they reflect what vendors and papers claim, not "
"audited live behaviour; I benchmarked no agent on real performance. It's a snapshot, not a film — 78 agents "
"frozen at one moment in a market that moves weekly. The tails are thin — one behavioral agent, basically no "
"black-swan ones, which is too few to generalize about the very corner that matters most. And the forced "
"single advantage that buys comparability also flattens hybrid tools. My stance is on the bottom line: none "
"of this invalidates the instrument, it scopes it. This is a transparent first map, built to be argued with — "
"which is exactly the empirical-validation agenda your Discussion section sets out. I see these as the "
"roadmap, not the verdict.")
footer(s, 13)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CONCLUSIONS / NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Conclusions & next steps")
title(s, "From a paper to a platform — and onward")

_, tf = box(s, MARGIN, Inches(2.0), Inches(6.1), Inches(3.3))
p = tf.paragraphs[0]; p.space_after = Pt(10)
run(p, "What stands today", 15, ACCENT, bold=True)
bl = tf
for txt in [
    "Your framework is executable, not just descriptive.",
    "A 78-agent benchmark with four queryable lenses.",
    "An empirical market map that confirms the paper's core thesis.",
    "Board-ready reporting for managers, developers, regulators.",
]:
    pp = bl.add_paragraph(); pp.space_after = Pt(7)
    run(pp, "✓  ", 14, ANA, bold=True); run(pp, txt, 13.5, TXT2)

px = Inches(7.0); pw = Inches(5.75)
rrect(s, px, Inches(2.0), pw, Inches(3.85), fill=SURFACE, line=BORDER, radius=0.04)
_, tf = box(s, px + Inches(0.34), Inches(2.22), pw - Inches(0.68), Inches(0.5))
run(tf.paragraphs[0], "WHERE IT GOES NEXT", 12, ACCENT, bold=True, spacing=1.8)
_, tf = box(s, px + Inches(0.34), Inches(2.78), pw - Inches(0.68), Inches(3.0))
nxts = [
    ("Validate", "inter-rater scoring + vendor confirmation of classifications."),
    ("Track over time", "snapshots → a longitudinal view of market drift."),
    ("Deepen the tails", "targeted hunt for behavioral & black-swan agents."),
    ("Tie to outcomes", "link classification to realized portfolio performance."),
    ("Generalize", "extend the framework beyond investment management."),
]
first = True
for head, body in nxts:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(9)
    run(p, "→  ", 14, ACCENT, bold=True)
    run(p, head + ".  ", 14, TXT, bold=True)
    run(p, body, 13, TXT2)
notes(s,
"To draw it together. What stands today, on the left: your framework is now executable rather than just "
"descriptive; there's a 78-agent benchmark with four lenses; an empirical market map that independently "
"confirms your central thesis; and reporting that reaches all three of your audiences. Where it goes next, "
"on the right, maps almost one-to-one onto your own future-research section: validate it with inter-rater "
"scoring and vendor confirmation; turn snapshots into a longitudinal study of how the market drifts; "
"deliberately go hunting in the thin tails — the behavioral and black-swan agents; connect classification to "
"realized portfolio outcomes; and finally test whether the framework generalizes beyond investment "
"management. The throughline of my internship was simple: take a strong idea on paper and prove it can live "
"as a working instrument. I think it can.")
footer(s, 14)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ══════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, Inches(0.07), fill=ACCENT)
if os.path.exists(logo):
    lw = Inches(2.2); lh = Inches(2.2 * 610 / 1562)
    s.shapes.add_picture(logo, MARGIN, Inches(0.7), lw, lh)

_, tf = box(s, MARGIN, Inches(2.9), Inches(11.9), Inches(1.4))
p = tf.paragraphs[0]; p.space_after = Pt(4)
run(p, "Thank you.", 40, TXT, bold=True)
p2 = tf.add_paragraph()
run(p2, "I'd value your critique of the classifications — you wrote the framework I built on.", 18, ACCENT)

rect(s, MARGIN, Inches(4.7), Inches(3.2), Pt(2), fill=ACCENT_DK)
_, tf = box(s, MARGIN, Inches(4.95), Inches(11.9), Inches(1.4))
for line, col in [
    ("Discussion — where would you challenge a classification?", TXT),
    ("Which agents are missing from the benchmark?", TXT2),
    ("Is the 'Full' autonomy tier a distinction worth keeping?", TXT2),
]:
    pp = tf.add_paragraph(); pp.space_after = Pt(6)
    run(pp, "·  ", 15, ACCENT, bold=True); run(pp, line, 15, col)
notes(s,
"Thank you. Genuinely — you wrote the framework, and getting to build on it has been the most engaging part "
"of the internship. I'd love this to be a discussion rather than just a presentation, so three openers: "
"Where would you challenge a specific classification? Which agents am I missing? And is the 'Full' autonomy "
"tier a distinction you'd keep, or did I over-engineer it? Over to you.")
footer(s, 15)


prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
